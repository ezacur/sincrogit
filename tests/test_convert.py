"""Readable .pptx extraction (convert.py) and its integration end to end:
committed .pptx -> markdown previews, worktree text, and history search."""

import io
import os

import pytest

pptx_lib = pytest.importorskip("pptx", reason="python-pptx not installed")

from pptx.util import Inches

from sincrogit.convert import pptx_bytes_to_md

from conftest import git, write


def _make_pptx(title="Resultados Q2", bullet="Ventas +12%", note="revisar cifra",
               cell="dato") -> bytes:
    """A real one-slide .pptx: title, a two-level bullet list, a table, notes."""
    prs = pptx_lib.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # title + content
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.text = bullet
    p2 = body.add_paragraph()
    p2.text = "detalle anidado"
    p2.level = 1
    table = slide.shapes.add_table(1, 2, Inches(1), Inches(4),
                                   Inches(4), Inches(1)).table
    table.cell(0, 0).text = cell
    table.cell(0, 1).text = "otro"
    slide.notes_slide.notes_text_frame.text = note
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_corrupt_pptx_returns_none_and_says_why(caplog):
    """Corrupt bytes and 'library missing' both yield None for callers, but a
    corrupt file must leave a log trace — '(binary or unavailable)' in a diff
    should be explainable."""
    import logging
    with caplog.at_level(logging.WARNING, logger="sincrogit.convert"):
        assert pptx_bytes_to_md(b"these bytes are not a zip archive") is None
    assert any("cannot read .pptx" in r.message for r in caplog.records)


def test_extracts_title_bullets_table_and_notes():
    md = pptx_bytes_to_md(_make_pptx())
    assert "## Slide 1 — Resultados Q2" in md
    assert "- Ventas +12%" in md
    assert "    - detalle anidado" in md      # bullet level -> indent
    assert "| dato | otro |" in md
    assert "> Notes: revisar cifra" in md


def test_conversion_is_deterministic():
    data = _make_pptx()
    assert pptx_bytes_to_md(data) == pptx_bytes_to_md(data)


def test_edge_cases():
    assert pptx_bytes_to_md(b"") == ""
    assert pptx_bytes_to_md(b"not a pptx at all") is None


def test_committed_pptx_reads_as_markdown(make_repo, make_engine):
    repo = make_repo({"deck.pptx": _make_pptx()})
    eng, _ = make_engine(repo, extra_includes=["**/*.pptx"])
    sha = git(repo, "rev-parse", "HEAD")
    text = eng.file_content_at("t", "deck.pptx", sha)  # raw path stays binary-ish
    assert eng.file_text_at("t", "deck.pptx", sha) != text
    assert "## Slide 1 — Resultados Q2" in eng.file_text_at("t", "deck.pptx", sha)
    assert "Ventas +12%" in eng.worktree_text("t", "deck.pptx")


def test_history_search_works_on_pptx(tmp_path, make_engine, monkeypatch):
    """Every commit gets an explicit, increasing committer date: version
    ordering (and blob collapsing) sorts by epoch, and same-second commits
    would tie arbitrarily — including the WIP the engine creates at setup."""
    import subprocess

    def commit_all(repo, msg, when):
        env = {**os.environ, "GIT_COMMITTER_DATE": when, "GIT_AUTHOR_DATE": when}
        subprocess.run(["git", "-C", repo, "add", "-A"], check=True,
                       capture_output=True)
        subprocess.run(["git", "-C", repo, "commit", "-m", msg], check=True,
                       capture_output=True, env=env)

    repo = str(tmp_path / "deck_repo")
    os.makedirs(repo)
    git(repo, "init", "-b", "main")
    git(repo, "config", "user.email", "t@example.com")
    git(repo, "config", "user.name", "T")
    write(repo, "deck.pptx", _make_pptx(bullet="Ventas +12%"))
    commit_all(repo, "feat: base", "2026-01-01T00:00:00")
    write(repo, "deck.pptx", _make_pptx(bullet="Ventas +20% y objetivo nuevo"))
    commit_all(repo, "feat: v2", "2026-01-01T01:00:00")

    # The engine's setup creates the WIP; date it too so it sorts newest.
    monkeypatch.setenv("GIT_COMMITTER_DATE", "2026-01-01T02:00:00")
    monkeypatch.setenv("GIT_AUTHOR_DATE", "2026-01-01T02:00:00")
    eng, _ = make_engine(repo, extra_includes=["**/*.pptx"])
    counts = [n for _s, n in eng.search_in_file_versions("t", "deck.pptx",
                                                         "objetivo nuevo")]
    # Newest first; the WIP shares v2's blob and collapses into it.
    assert counts == [1, 0]


def test_gitattributes_mapped_for_pptx(make_repo, make_engine):
    repo = make_repo({"deck.pptx": _make_pptx()})
    make_engine(repo, extra_includes=["**/*.pptx"])
    attrs = open(os.path.join(repo, ".gitattributes"), encoding="utf-8").read()
    assert "*.pptx -text" in attrs
