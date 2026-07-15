"""Readable-text extraction for binary documents pandoc can't read (.pptx).

A .pptx is converted to markdown IN-PROCESS via python-pptx: slide titles,
text frames (with bullet indent levels), tables and speaker notes — the text
content a diff actually cares about. Layout, images and animations are
intentionally ignored: they don't diff as text. The output is deterministic
(slide order, then shape order as stored in the XML), which is what makes
version-to-version diffs meaningful.

python-pptx is OPTIONAL, like pandoc: without it .pptx files are still
versioned, just as opaque blobs (no readable preview/diff). Unlike the .docx
path there is no git textconv driver here (that needs an external executable),
so a .pptx snapshot is gated by BYTES, not by content: any resave counts.
"""

import io
import logging

log = logging.getLogger("sincrogit.convert")

_pptx_checked = False
_pptx_ok = False


def pptx_available() -> bool:
    """Is python-pptx importable? Probed once, cached — mirrors resolve_pandoc:
    repos that never touch a .pptx never pay the import."""
    global _pptx_checked, _pptx_ok
    if not _pptx_checked:
        try:
            import pptx  # noqa: F401
            _pptx_ok = True
            log.info("python-pptx found: .pptx previews/diffs will be readable")
        except Exception:  # noqa: BLE001 — a broken install counts as absent
            _pptx_ok = False
        _pptx_checked = True
    return _pptx_ok


def _shape_lines(shape) -> list:
    """Markdown lines for one shape (recursing into groups). Defensive: an odd
    shape (chart, SmartArt, foreign object) yields nothing rather than raising."""
    lines = []
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                lines += _shape_lines(sub)
            return lines
        if getattr(shape, "has_table", False):
            for row in shape.table.rows:
                cells = [" ".join(cell.text.split()) for cell in row.cells]
                lines.append("| " + " | ".join(cells) + " |")
            return lines
        if getattr(shape, "has_text_frame", False):
            for para in shape.text_frame.paragraphs:
                text = " ".join(para.text.split())
                if text:
                    lines.append("    " * para.level + "- " + text)
    except Exception:  # noqa: BLE001 — skip the shape, keep the slide
        pass
    return lines


def pptx_bytes_to_md(data: bytes) -> str | None:
    """Convert .pptx bytes to markdown. '' for empty input, None when
    python-pptx is missing or the file isn't a readable presentation."""
    if not data:
        return ""
    if not pptx_available():
        return None
    from pptx import Presentation
    try:
        prs = Presentation(io.BytesIO(data))
    except Exception as e:  # noqa: BLE001 — corrupt/foreign bytes: not convertible
        # Same None as "library missing" for callers, but a distinguishable
        # trace for the user: "(binary or unavailable)" in a diff should be
        # explainable — a truncated/corrupt .pptx is worth knowing about.
        log.warning("cannot read .pptx (%d bytes): %s: %s",
                    len(data), type(e).__name__, e)
        return None
    out = []
    for num, slide in enumerate(prs.slides, 1):
        title, title_id = "", None
        try:
            t = slide.shapes.title
            if t is not None:
                title_id = t.shape_id
                title = " ".join(t.text.split())
        except Exception:  # noqa: BLE001 — layouts without a title placeholder
            pass
        out.append(f"## Slide {num}" + (f" — {title}" if title else ""))
        for shape in slide.shapes:
            if title_id is not None and getattr(shape, "shape_id", None) == title_id:
                continue  # the title is already in the heading
            out += _shape_lines(shape)
        try:
            notes = (slide.notes_slide.notes_text_frame.text.strip()
                     if slide.has_notes_slide else "")
        except Exception:  # noqa: BLE001
            notes = ""
        out += [f"> Notes: {ln.strip()}" for ln in notes.splitlines() if ln.strip()]
        out.append("")  # blank line between slides
    return "\n".join(out).rstrip() + "\n"
